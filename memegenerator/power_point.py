"""Generate PowerPoint file"""
from pptx import Presentation
from pptx.util import Inches, Pt
from os.path import splitext


def gen_ppt(img_path, mytxt):
    """Generte ppt with picture."""
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "This world is nothing without art."
    subtitle.text = "if there is not Art, there is no Life."

    ############### Picture Slide #########

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    top = Inches(1)
    left = Inches(1)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    newpath_tuple = splitext(img_path)
    new_path = newpath_tuple[0] + '.ppt'

    ############### Slides     ############
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "The Picture on previous slide is captioned with quotation. \nThe same text is used on this slide as well. See below \n" + mytxt

    p = tf.add_paragraph()
    p.text = mytxt + '\n\n\nThis text can be parsed in any way, for example can be reversed: \n' + mytxt[::-1] 
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = mytxt
    p.font.size = Pt(40)


    prs.save(new_path)
    return new_path
